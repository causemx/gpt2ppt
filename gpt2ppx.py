import os
import openai
import requests 
import json
import extension
from argparse import ArgumentParser as ap
from pptx import Presentation as pt


openai.api_key = extension.OPENAI_API_KEY

url='https://api.openai.com/v1/chat/completions'

class GenerationError(Exception):
    def __init__(self, f, *args):
        super().__init__(args)
        self.f = f

    def __str__(self):
        return f'There are something wrong when generating'

def analyze_res(data):
    # byte to string
    data_str = data.decode('utf-8')
    json_data = json.loads(data_str)
    response = json_data['choices'][0]['message']['content']
    return response

def add_slide(prs, layout, title, content):
    """Return slide newly added to `prs` using `layout` and having `title`."""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content
    return slide

def generate(args):
    prs = pt()
    input_file = args.input_file
    output_file = args.output_file
    output = "./output.pptx" if not output_file else output_file

    title_slide_layout = prs.slide_layouts[1]
    # IMPORTANT: set slogan for starting
    page_index = 1
    prompt = ''
    print("Hello")
    try:
        with open(input_file) as topo_file:
            for line in topo_file:
                print(f"generating slide page: {page_index}")
                question, topic = line.split(',')
                # prompt += input(":")+'\n'
                prompt += question + '\n'
                payload = {
                    "model": "gpt-3.5-turbo",
                    "messages": [{"role": "user", "content": prompt}]
                    #"n": 10,
                    #"temperature: 0.7
                }
                headers = {
                    "Authorization": f"Bearer {openai.api_key}",
                    "Content-Type":"application/json"
                } 
                r = requests.post(url, data=json.dumps(payload), headers=headers)
                try:
                    res = analyze_res(r.content)   
                    prompt +=  res+'\n'
                    add_slide(prs, title_slide_layout, topic, res)
                except GenerationError:
                    print("error:",r.content)
                    return 1
                
                page_index = page_index + 1
                
        prs.save(output)       
    except FileNotFoundError as fe:
        print("error:", str(fe))
        return 1

    return 0

if __name__ == "__main__":
    parser = ap(description="gpt2ppt commands")
    parser.add_argument("input_file", help=".txt, format:[question],[topic]")
    parser.add_argument("-o", "--output_file", help="output destination path (folder)")
    args = parser.parse_args()
    error_code = generate(args)
    if error_code != 0:
        print("exit...")
        exit(error_code)
    

    