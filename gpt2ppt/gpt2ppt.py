import os
import openai
import json
from gpt2ppt import extension
from .parse import Parser
from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

openai.api_key = os.getenv("OPENAI_API_KEY")

MODEL_ENGINE = "gpt-3.5-turbo"


class GenerationError(Exception):
    def __init__(self, f, *args):
        super().__init__(args)
        self.f = f

    def __str__(self):
        return 'There are something wrong when generating'


class Slide:
    _instance = None
    SLD_LAYOUT_TITLE_AND_CONTENT = 1

    def __new__(cls):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
        return cls._instance
    
    def __init__(self, input=None, output=None) -> None:
        self.prs = Presentation()
    
    def add_slide(self):
        self.slide = self.prs.slides.add_slide(
            self.prs.slide_layouts[self.SLD_LAYOUT_TITLE_AND_CONTENT])
        self.shapes = self.slide.shapes
    
    def add_title(self, _title):
        title_shape = self.shapes.title
        title_shape.text = _title
    
    def add_body(self, _body, level=0):
        body_shape = self.shapes.placeholders[1]
        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        # tf.text = _body

        try:
            k, v = _body.split(':')
            if v.startswith("@prompt"):
                prompt = v.split("@prompt",1)[1].lstrip()
                p = tf.add_paragraph()
                # p.text = "{} - {}".format(k, gen(prompt))
                p.text = "{} - {}".format(k, prompt)
                p.level = level-1
                
                font = p.font
                font.name = 'Calibri'
                font.size = Pt(16)
                font.bold = True
                font.italic = True  # cause value to be inherited from theme
                font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
                
            else:
                raise Exception   
        except:
            p = tf.add_paragraph()
            p.text = _body
            p.level = level-1

    def save(self, output_file):
        self.prs.save(output_file)

def generate(args):
    try:
        with open(args.input_file) as topo_file:
            # Setup presentation
            s = Slide()
            
            doc = json.load(topo_file)
            p = Parser()
            p.parse(doc)
            _queue = p.q
            while not _queue.empty():
                subline = _queue.get()
                # Check title
                if subline[0].startswith("@title"):
                    s.add_slide()
                    s.add_title(subline[0].split("@title",1)[1].lstrip())
                else:
                    s.add_body(subline[0], subline[1])
            s.save(args.outputFile)
    except FileNotFoundError as fe:
        print("error:", str(fe))
        return 1
    return 0

def gen(prompts, max=30):
    for prompt in prompts:
        response = openai.ChatCompletion.create(
            model=MODEL_ENGINE, 
            messages=[{"role": "user", 
                       "content": "{} in {} words".format(prompt, max)}])   
    return response['choices'][0]['message']['content']