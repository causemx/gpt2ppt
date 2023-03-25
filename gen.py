import os
import openai
import requests 
import json
import extension
from pptx import Presentation as pt


openai.api_key = extension.OPENAI_API_KEY

url='https://api.openai.com/v1/chat/completions'

class GenerationError(Exception):
    def __init__(self, f, *args):
        super().__init__(args)
        self.f = f

    def __str__(self):
        return f'There are something wrong when generating'

def analyze_res(data):
    # byte to string
    data_str = data.decode('utf-8')
    json_data = json.loads(data_str)
    response = json_data['choices'][0]['message']['content']
    return response

def add_slide(prs, layout, title, content):
    """Return slide newly added to `prs` using `layout` and having `title`."""
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content
    return slide

if __name__ == "__main__":

    prs = pt()

    title_slide_layout = prs.slide_layouts[1]
    """slides = [
        add_slide(prs, title_slide_layout, "Summary Table", "this is summary"),
        add_slide(prs, title_slide_layout, "New Table", "this is new table"),
        add_slide(prs, title_slide_layout, "Old Table", "this is old table"),
        ]
    prs.save('output.pptx')"""
   


    # IMPORTANT: set slogan for starting
    prompt = ''
    print("Hello")

    with open('input.txt') as topo_file:
        for line in topo_file:
            question, topic = line.split(',')
            
            # prompt += input(":")+'\n'
            prompt += question + '\n'
            payload = {
                "model": "gpt-3.5-turbo",
                "messages": [{"role": "user", "content": prompt}]
                #"n": 10,
                #"temperature: 0.7
            }

            headers = {
                "Authorization": f"Bearer {openai.api_key}",
                "Content-Type":"application/json"
            } 
            r = requests.post(url, data=json.dumps(payload), headers=headers)
            try:
                res = analyze_res(r.content)   
                prompt +=  res+'\n'
                add_slide(
                    prs, 
                    title_slide_layout, 
                    question, 
                    res)

            except GenerationError:
                print("error:",r.content)
                break 
    prs.save('output.pptx') 

    """while True:
        prompt += input(":")+'\n'
        payload = {
            "model": "gpt-3.5-turbo",
            "messages": [{"role": "user", "content": prompt}]
            #"n": 10,
            #"temperature: 0.7
        }

        headers = {
            "Authorization": f"Bearer {openai.api_key}",
            "Content-Type":"application/json"
        }
        
        r = requests.post(url, data=json.dumps(payload), headers=headers)
        try:
            res=analyze_res(r.content)   
            prompt +=  res+'\n'
            print(res)
        except:
            print("error:",r.content)"""
